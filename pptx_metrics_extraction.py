from pptx_metric_mapping import process_json_file
from fastapi import FastAPI, UploadFile, File, APIRouter
from fastapi.responses import JSONResponse
import os, re, json
import zipfile
from pptx import Presentation
import pandas as pd
import pytesseract
from PIL import Image
import xml.etree.ElementTree as ET
from openai import OpenAI
from dotenv import load_dotenv
import logging

# Set up logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

dotenv_path = os.path.expanduser("~/.env")

# Load the environment variables from the specified file
load_dotenv(dotenv_path)
client = OpenAI()
app = FastAPI()
pptx_router = APIRouter()
# pytesseract.pytesseract.tesseract_cmd = r"C:\Program Files\Tesseract-OCR\tesseract.exe"

async def process_pptx(file: UploadFile = File(...)):
    # Process the uploaded PPTX file
    result = await process_pptx_content(file)
    return {"result": result}

def extract_text(slide):
    """Extract text from a slide"""
    text = ""
    for shape in slide.shapes:
        if hasattr(shape, "text"):
            text += shape.text + "\n"
    return text

def extract_tables(slide):
    """Extract tables from a slide"""
    tables = []
    for shape in slide.shapes:
        if shape.has_table:
            table_data = []
            table = shape.table
            for row in table.rows:
                row_data = [cell.text for cell in row.cells]
                table_data.append(row_data)
            tables.append(table_data)
    return tables

def get_slide_image_mappings(slides_rels_dir):
    """Get mappings between slides and their images"""
    slide_images = {}
    
    for rel_file in os.listdir(slides_rels_dir):
        if rel_file.endswith('.rels'):
            slide_num = rel_file.split('slide')[1].split('.xml')[0]
            rel_path = os.path.join(slides_rels_dir, rel_file)
            
            tree = ET.parse(rel_path)
            root = tree.getroot()
            
            images = []
            for rel in root.findall('{http://schemas.openxmlformats.org/package/2006/relationships}Relationship'):
                target = rel.get('Target')
                if target and '../media/' in target and any(
                    target.lower().endswith(ext) for ext in ['.png', '.jpg', '.jpeg', '.bmp']):
                    image_file = os.path.basename(target)
                    images.append(image_file)
            
            if images:
                slide_images[slide_num] = images
    
    return slide_images

def get_chart_excel_mappings(slides_rels_dir, charts_rels_dir):
    """Get mappings between slides and their Excel files"""
    slide_to_excel = {}
    
    for slide_rel_file in os.listdir(slides_rels_dir):
        if slide_rel_file.endswith('.rels'):
            slide_number = slide_rel_file.split('slide')[1].split('.xml')[0]
            slide_rel_path = os.path.join(slides_rels_dir, slide_rel_file)
            
            slide_tree = ET.parse(slide_rel_path)
            slide_root = slide_tree.getroot()
            
            for rel in slide_root.findall('{http://schemas.openxmlformats.org/package/2006/relationships}Relationship'):
                target = rel.get('Target')
                if target and target.startswith("../charts/chart") and target.endswith(".xml"):
                    chart_file = os.path.basename(target)
                    chart_number = chart_file.split('chart')[1].split('.xml')[0]
                    
                    chart_rel_file = f"chart{chart_number}.xml.rels"
                    chart_rel_path = os.path.join(charts_rels_dir, chart_rel_file)
                    
                    if os.path.exists(chart_rel_path):
                        chart_tree = ET.parse(chart_rel_path)
                        chart_root = chart_tree.getroot()
                        
                        for chart_rel in chart_root.findall('{http://schemas.openxmlformats.org/package/2006/relationships}Relationship'):
                            chart_target = chart_rel.get('Target')
                            if chart_target and chart_target.startswith("../embeddings/") and chart_target.endswith(".xlsx"):
                                embedded_excel = os.path.basename(chart_target)

                                if slide_number not in slide_to_excel:
                                    slide_to_excel[slide_number] = []
                                slide_to_excel[slide_number].append(embedded_excel)
    
    return slide_to_excel

def ocr_image(image_path):
    """Perform OCR on an image"""
    try:
        image = Image.open(image_path)
        text = pytesseract.image_to_string(image)
        return text.strip()
    except Exception as e:
        print(f"Error performing OCR on {image_path}: {str(e)}")
        return ""

from dotenv import load_dotenv
load_dotenv() 

def summarize_content(content):
    completion = client.chat.completions.create(
        model="gpt-4o-mini",
        messages=[
            {"role": "system", "content": "You are a helpful assistant."},
            {
                "role": "user",
                "content": f"Summarize the following content:\n\n{content}"
            }
        ], temperature= 0.0
    )
    summary = completion.choices[0].message.content.strip()
    return summary

def extract_slide_content(content, slide_number):

    slide_pattern = fr"# Slide {slide_number}\n\n(.*?)(?=# Slide \d+|$)"
    match = re.search(slide_pattern, content, re.DOTALL)
    
    if match:
        return match.group(1).strip()
    return ""


from pptx_prompt import prompt as pmt
def generate_slide_json(content, slide_number):
    """
    Generate JSON for a specific slide
    
    Args:
        content (str): Combined markdown content
        slide_number (int): Slide number to process
        prompt (str): Extraction prompt
    
    Returns:
        dict: JSON representation of the slide
    """
    slide_content = extract_slide_content(content, slide_number)
    
    if not slide_content:
        return None
    
    completion = client.chat.completions.create(
        model="gpt-4o-mini",
        messages=[
            {"role": "system", "content": "You are an expert on social events who is given markdown files and creates valid json files from given data."},
            {
                "role": "user", 
                "content": f"{pmt}:\n\n{slide_content}"
            }
        ]
    )
    
    summary_raw = completion.choices[0].message.content.strip()
    summary_json = re.sub(r'```json|```', '', summary_raw).strip()
    
    try:
        return json.loads(summary_json)
    except json.JSONDecodeError:
        print(f"Error decoding JSON for slide {slide_number}")
        print(f"RAW SUMMARY: {summary_raw}")
        return None

def process_pptx_to_json(presentation, markdowns_dir):
    # Combined markdown file path
    combined_markdown_path = os.path.join(markdowns_dir, "combined_slides.md")
    
    # Read combined markdown
    with open(combined_markdown_path, "r", encoding='utf-8') as combined_md_file:
        combined_content = combined_md_file.read()
    
    # Store individual slide JSONs
    slide_jsons = {}
    
    # Process each slide
    for slide_number in range(1, len(presentation.slides) + 1):
        slide_json = generate_slide_json(combined_content, slide_number)
        
        if slide_json:
            # Save individual slide JSON
            slide_json_path = os.path.join(markdowns_dir, f"slide_{slide_number}_metrics.json")
            with open(slide_json_path, "w", encoding='utf-8') as json_file:
                json.dump(slide_json, json_file, indent=2)
            
            slide_jsons[f"slide_{slide_number}"] = slide_json
            print(f"JSON for Slide {slide_number} saved to {slide_json_path}")
    
    # Combine all slide JSONs
    combined_slides_json = {
        "slides": slide_jsons
    }
    
    # Save combined slides JSON
    combined_json_path = os.path.join(markdowns_dir, "combined_slides_metrics.json")
    with open(combined_json_path, "w", encoding='utf-8') as combined_json_file:
        json.dump(combined_slides_json, combined_json_file, indent=2)
    
    print(f"Combined slides JSON saved to {combined_json_path}")

@app.post("/process_pptx/")
async def process_pptx_content(pptx_file: UploadFile = File(...)):
    # Save the uploaded file temporarily
    try:
        os.makedirs("/tmp", exist_ok=True)
        temp_file_path = f"/tmp/{pptx_file.filename}"
        with open(temp_file_path, "wb") as buffer:
            buffer.write(await pptx_file.read())

        # Process the PPTX file
        extracted_dir = os.path.join(r"../extracted_content", "everything", pptx_file.filename.split('.')[0])
        markdowns_dir = os.path.join(extracted_dir, "markdowns")
        
        # Create necessary directories
        os.makedirs(markdowns_dir, exist_ok=True)
        
        # Extract PPTX contents
        print(f"Extracting PPTX contents to {extracted_dir}...")
        with zipfile.ZipFile(temp_file_path, 'r') as pptx:
            pptx.extractall(extracted_dir)
        
        # Define important directories
        slides_rels_dir = os.path.join(extracted_dir, 'ppt', 'slides', '_rels')
        charts_rels_dir = os.path.join(extracted_dir, 'ppt', 'charts', '_rels')
        embeddings_dir = os.path.join(extracted_dir, 'ppt', 'embeddings')
        media_dir = os.path.join(extracted_dir, 'ppt', 'media')
        
        # Get slide-to-image and slide-to-excel mappings
        slide_images = get_slide_image_mappings(slides_rels_dir)
        slide_to_excel = get_chart_excel_mappings(slides_rels_dir, charts_rels_dir)
        
        # Load presentation
        presentation = Presentation(temp_file_path)
        
        # Process each slide
        for i, slide in enumerate(presentation.slides):
            slide_number = i + 1
            markdown_file_path = os.path.join(markdowns_dir, f"slide_{slide_number}.md")
            
            print(f"\nProcessing Slide {slide_number}")
            
            with open(markdown_file_path, "w", encoding='utf-8') as md_file:
                # Write slide number
                md_file.write(f"# Slide {slide_number}\n\n")
                
                # Extract and write text content
                text_content = extract_text(slide)
                md_file.write(f"## Text Content\n\n{text_content}\n\n")
                
                # Extract and write table content
                tables = extract_tables(slide)
                if tables:
                    md_file.write("## Tables\n\n")
                    for table_index, table in enumerate(tables):
                        md_file.write(f"### Table {table_index + 1}\n\n")
                        for row in table:
                            md_file.write("| " + " | ".join(row) + " |\n")
                        md_file.write("\n")
                
                # Process images
                if str(slide_number) in slide_images:
                    md_file.write("## Images\n\n")
                    for img_file in slide_images[str(slide_number)]:
                        img_path = os.path.join(media_dir, img_file)
                        print(f"  Processing image: {img_file}")
                        
                        # Perform OCR
                        ocr_text = ocr_image(img_path)
                        
                        md_file.write(f"### Image: {img_file}\n\n")
                        if ocr_text:
                            print(f"  OCR Text found in {img_file}")
                            md_file.write(f"OCR Text:\n```\n{ocr_text}\n```\n\n")
                        else:
                            print(f"  No text found in {img_file}")
                            md_file.write("No text found in image\n\n")
                
                if str(slide_number) in slide_to_excel:
                    md_file.write("## Excel Data\n\n")
                    for excel_file in slide_to_excel[str(slide_number)]:
                        excel_path = os.path.join(embeddings_dir, excel_file)
                        if os.path.exists(excel_path):
                            print(f"  Processing Excel file: {excel_file}")
                            try:
                                df = pd.read_excel(excel_path)
                                md_file.write(f"### {excel_file}\n\n")
                                md_file.write(df.to_markdown(index=False))
                                md_file.write("\n\n")
                            except Exception as e:
                                print(f"  Error processing Excel file {excel_file}: {str(e)}")
                                md_file.write(f"Error processing Excel file: {excel_file}: {str(e)} \n\n")
            
            print(f"  Saved content to {markdown_file_path}")
        # Combine all markdown files into a single file
        combined_markdown_path = os.path.join(markdowns_dir, "combined_slides.md")
        with open(combined_markdown_path, "w", encoding='utf-8') as combined_md_file:
            for slide_number in range(1, len(presentation.slides) + 1):
                markdown_file_path = os.path.join(markdowns_dir, f"slide_{slide_number}.md")
                with open(markdown_file_path, "r", encoding='utf-8') as md_file:
                    combined_md_file.write(md_file.read())
                    combined_md_file.write("\n\n")  # Add some space between slides

        print(f"Combined markdown file saved to {combined_markdown_path}")
        with open(combined_markdown_path, "r", encoding='utf-8') as combined_md_file:
            combined_content = combined_md_file.read()

        # Send the content to the LLM for summarization
        summary = summarize_content(combined_content)

        # Save the summary to a markdown file
        summary_file_path = os.path.join(markdowns_dir, "summary.md")
        try:
            with open(summary_file_path, "w", encoding='utf-8') as summary_file:
                summary_file.write("# Summary\n\n")
                summary_file.write(summary)
            print(f"Summary saved to {summary_file_path}")
        except Exception as e:
            print(f"Error saving summary: {str(e)}")
        process_pptx_to_json(presentation, markdowns_dir)
        # Define input and output file paths
        input_json_path = os.path.join(markdowns_dir, "combined_slides_metrics.json")
        output_metrics_path = os.path.join(markdowns_dir, "processed_metrics.json")

        logger.info(f"Attempting to start metric mapping")
        logger.info(f"Input JSON path: {input_json_path}")
        logger.info(f"Output metrics path: {output_metrics_path}")

        try:
            # Directly await the async function instead of using create_task
            await process_json_file(input_json_path, output_metrics_path)
            logger.info(f"Metric mapping completed successfully")
        except Exception as e:
            logger.error(f"Error in metric mapping: {str(e)}")
            logger.error(f"Full traceback:", exc_info=True)

        # Remove the temporary file
        os.remove(temp_file_path)

        return JSONResponse(content={
            "message": "PPTX processed successfully", 
            "metric_mapping_status": "Processing started in background"
        })

    except Exception as e:
        return JSONResponse(
            status_code=500, 
            content={"message": f"Error processing PPTX: {str(e)}"}
        )

async def run_metric_mapping(input_file: str, output_file: str):
    try:
        await process_json_file(input_file, output_file)
        logger.info(f"Metric mapping completed. Output written to {output_file}")
    except Exception as e:
        logger.error(f"Error in metric mapping: {str(e)}")

@app.get("/metric_mapping_status")
async def get_metric_mapping_status():
    """
    Endpoint to check the status of metric mapping process
    """
    try:
        # Check if processed_metrics.json exists
        if os.path.exists("processed_metrics.json"):
            return JSONResponse(content={
                "status": "completed",
                "file": "processed_metrics.json"
            })
        else:
            return JSONResponse(content={
                "status": "processing",
                "message": "Metric mapping is in progress"
            })
    except Exception as e:
        return JSONResponse(
            status_code=500, 
            content={"message": f"Error checking status: {str(e)}"}
        )